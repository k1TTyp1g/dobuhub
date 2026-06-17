"""
DocuHub 文档管家 — 最终版
左：文件夹 + 最近打开 | 右上：文件列表 | 右下：文件预览
"""
import tkinter as tk
from tkinter import ttk, messagebox, filedialog
from ttkthemes import ThemedTk
import threading
import os
import sys
import json
import ctypes
from ctypes import wintypes
from datetime import datetime
from PIL import Image, ImageTk

COLOR_PRIMARY = "#2B579A"
COLOR_ACCENT = "#1A73E8"
COLOR_BG = "#F5F5F5"
COLOR_CARD = "#FFFFFF"
FONT_FAMILY = "Microsoft YaHei"
COLOR_DOC = "#E8F5E9"

ALLOWED_EXTENSIONS = {
    ".doc", ".docx", ".txt", ".pdf", ".rtf", ".wps", ".odt",
    ".xls", ".xlsx", ".csv", ".et",
    ".ppt", ".pptx", ".pps", ".dps",
    ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp",
    ".zip", ".rar", ".7z", ".eml", ".msg",
}

BLOCKED_EXTENSIONS = {
    ".lnk", ".bat", ".cmd", ".exe", ".dll", ".sys", ".tmp",
    ".ini", ".log", ".dat", ".cache", ".db", ".idx",
    ".msi", ".com", ".scr", ".pif", ".vbs", ".ps1",
    ".o", ".obj", ".lib", ".class",
}

SKIP_DIRS = {
    "windows", "winnt", "program files", "program files (x86)",
    "programdata", "system volume information",
    "boot", "recovery", "python", "python38", "node_modules",
    ".git", ".svn", "__pycache__", "msocache",
    "perflogs", "config.msi", "tmp", "temp",
}

SKIP_DIRS_LOWER = {d.lower() for d in SKIP_DIRS}

TEXT_EXTENSIONS = {".txt", ".md", ".py", ".js", ".html", ".css", ".json", ".xml", ".csv", ".log", ".ini", ".yaml", ".yml"}

TYPE_LABELS = {
    ".doc": "Word", ".docx": "Word", ".wps": "WPS",
    ".ppt": "PPT", ".pptx": "PPT", ".xls": "Excel", ".xlsx": "Excel",
    ".txt": "文本", ".pdf": "PDF",
    ".png": "图片", ".jpg": "图片", ".jpeg": "图片", ".gif": "图片",
    ".bmp": "图片", ".webp": "图片", ".tiff": "图片",
    ".zip": "压缩包", ".rar": "压缩包", ".7z": "压缩包",
    ".eml": "邮件", ".msg": "邮件", ".csv": "CSV", ".et": "WPS表格",
    ".rtf": "富文本", ".odt": "ODT", ".pps": "PPT", ".dps": "WPS演示",
}

# Windows Shell icon API
class SHFILEINFOW(ctypes.Structure):
    _fields_ = [
        ('hIcon', ctypes.c_void_p),
        ('iIcon', ctypes.c_int),
        ('dwAttributes', ctypes.c_ulong),
        ('szDisplayName', ctypes.c_wchar * 260),
        ('szTypeName', ctypes.c_wchar * 80),
    ]

SHGFI_ICON = 0x100
SHGFI_SMALLICON = 0x1
SHGFI_LARGEICON = 0x0
SHGFI_USEFILEATTRIBUTES = 0x10


def should_show_file(filename):
    if filename.startswith(".") or filename.startswith("~"):
        return False
    if filename.lower() in ("desktop.ini", "thumbs.db"):
        return False
    ext = os.path.splitext(filename)[1].lower()
    return ext not in BLOCKED_EXTENSIONS and ext in ALLOWED_EXTENSIONS


def get_type_label(ext):
    return TYPE_LABELS.get(ext, ext.strip(".").upper() if ext else "文件")


def format_size(size_bytes):
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 ** 2:
        return f"{size_bytes / 1024:.1f} KB"
    elif size_bytes < 1024 ** 3:
        return f"{size_bytes / (1024 ** 2):.1f} MB"
    else:
        return f"{size_bytes / (1024 ** 3):.2f} GB"


def send2trash(path):
    path = os.path.abspath(path)
    if not os.path.isfile(path):
        return
    class SHFILEOPSTRUCT(ctypes.Structure):
        _fields_ = [
            ('hwnd', ctypes.c_void_p), ('wFunc', wintypes.UINT),
            ('pFrom', ctypes.c_wchar_p), ('pTo', ctypes.c_wchar_p),
            ('fFlags', wintypes.UINT), ('fAnyOperationsAborted', ctypes.c_bool),
            ('hNameMappings', ctypes.c_void_p), ('lpszProgressTitle', ctypes.c_wchar_p),
        ]
    buf = ctypes.create_unicode_buffer(path, len(path) + 2)
    ctypes.windll.shell32.SHFileOperationW(ctypes.byref(SHFILEOPSTRUCT(
        hwnd=None, wFunc=3, pFrom=ctypes.cast(buf, ctypes.c_wchar_p),
        pTo=None, fFlags=0x40|0x10|0x04|0x0400,
        fAnyOperationsAborted=False, hNameMappings=None, lpszProgressTitle=None)))


class DocuHubApp:
    def __init__(self):
        # 使用 ttkthemes 现代主题
        self.root = ThemedTk(theme="breeze")
        self.root.title("文档管家 DocuHub")
        self.root.geometry("1200x750")
        self.root.minsize(700, 500)

        self.current_dir = ""
        self.folder_list = []
        self.saved_folders = []
        self.locked_folders = set()
        self._file_data = []
        self._file_path_map = {}
        self._sort_col = None
        self._sort_rev = False
        self._type_filter = set()
        self._icon_cache = {}
        self._filter_vars = {}
        self.recent_files = []
        self.RECENT_PATH = os.path.join(os.path.dirname(os.path.abspath(sys.argv[0])), "recent_files.json")

        self._build_layout()
        self._load_folders()
        self._load_recent()

    def run(self):
        self.root.mainloop()

    def _build_layout(self):
        main = tk.Frame(self.root, bg=COLOR_BG)
        main.pack(fill=tk.BOTH, expand=True, padx=4, pady=4)

        # ── 左面板：文件夹 + 最近打开 ──
        left = tk.Frame(main, bg=COLOR_CARD, bd=1)
        left.pack(side=tk.LEFT, fill=tk.Y, padx=(0, 2))
        left.config(width=200)
        left.pack_propagate(False)

        tk.Label(left, text="文件夹", font=(FONT_FAMILY, 11, "bold"),
                 bg=COLOR_CARD).pack(pady=(8, 0), padx=8, anchor=tk.W)
        sf = tk.Frame(left, bg=COLOR_CARD)
        sf.pack(fill=tk.X, padx=6, pady=2)
        self.global_search_var = tk.StringVar()
        gse = tk.Entry(sf, textvariable=self.global_search_var,
                        font=(FONT_FAMILY, 9))
        gse.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 2))
        gse.bind("<Return>", lambda e: self._start_global_search())
        ttk.Button(sf, text="搜",
                  command=self._start_global_search,
                  width=3).pack(side=tk.RIGHT)

        self.folder_listbox = tk.Listbox(left, font=(FONT_FAMILY, 9), selectbackground=COLOR_ACCENT, highlightthickness=0, bd=0,
                                          height=8)
        self.folder_listbox.pack(fill=tk.X, padx=4, pady=2)
        self.folder_listbox.bind("<<ListboxSelect>>", self._on_folder_select)
        self.folder_listbox.bind("<Button-3>", self._on_folder_right_click)

        self.folder_menu = tk.Menu(self.root, tearoff=0, font=(FONT_FAMILY, 9))
        self.folder_menu.add_command(label="刷新", command=self._folder_refresh)
        self.folder_menu.add_command(label="在资源管理器中打开", command=self._folder_explorer)
        self.folder_menu.add_separator()
        self.folder_menu.add_command(label="锁定", command=self._folder_lock)
        self.folder_menu.add_command(label="解锁", command=self._folder_unlock)
        self.folder_menu.add_separator()
        self.folder_menu.add_command(label="从列表删除", command=self._folder_remove)

        bf = tk.Frame(left, bg=COLOR_CARD)
        bf.pack(fill=tk.X, padx=6, pady=2)
        self.scan_btn = ttk.Button(bf, text="全盘扫描",
                                   command=self._start_full_scan)
        self.scan_btn.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 2))
        ttk.Button(bf, text="添加",
                  command=self._add_folder).pack(side=tk.RIGHT, padx=(2, 0))

        # 最近打开
        tk.Label(left, text="最近打开", font=(FONT_FAMILY, 10, "bold"),
                 bg=COLOR_CARD).pack(pady=(6, 0), padx=8, anchor=tk.W)
        self.recent_listbox = tk.Listbox(left, font=(FONT_FAMILY, 9), selectbackground=COLOR_ACCENT, highlightthickness=0, bd=0,
                                          height=6)
        self.recent_listbox.pack(fill=tk.X, padx=4, pady=2)
        self.recent_listbox.bind("<Double-1>", self._open_recent)
        self.recent_listbox.bind("<Button-3>", self._recent_right_click)

        # 最近打开右键菜单
        self.recent_menu = tk.Menu(self.root, tearoff=0, font=(FONT_FAMILY, 9))
        self.recent_menu.add_command(label="打开", command=self._recent_open)
        self.recent_menu.add_command(label="打开所在文件夹", command=self._recent_open_location)
        self.recent_menu.add_separator()
        self.recent_menu.add_command(label="从列表删除", command=self._recent_remove)

        # ── 右侧：文件列表 + 预览（可拖动分割） ──
        right = tk.Frame(main, bg=COLOR_BG)
        right.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=2)

        paned = ttk.PanedWindow(right, orient=tk.VERTICAL)
        paned.pack(fill=tk.BOTH, expand=True)

        # 文件列表
        list_frame = tk.Frame(paned, bg=COLOR_CARD, bd=1)
        paned.add(list_frame, weight=3)

        top = tk.Frame(list_frame, bg=COLOR_CARD)
        top.pack(fill=tk.X, padx=6, pady=(6, 2))
        self.dir_label = tk.Label(top, text="(选择文件夹或搜索)",
                                  font=(FONT_FAMILY, 9), bg=COLOR_CARD, anchor=tk.W)
        self.dir_label.pack(side=tk.LEFT, fill=tk.X, expand=True)
        ttk.Button(top, text="刷新",
                  command=self._refresh_file_list).pack(side=tk.RIGHT, padx=2)
        tk.Label(top, text="搜索:", bg=COLOR_CARD).pack(side=tk.RIGHT, padx=2)
        self.search_entry = tk.Entry(top, font=(FONT_FAMILY, 9), width=15)
        self.search_entry.pack(side=tk.RIGHT, padx=2)
        self.search_entry.bind("<KeyRelease>", lambda e: self._debounce_search())

        # 类型筛选栏（内置勾选框）
        self.filter_bar = tk.Frame(list_frame, bg=COLOR_CARD)
        self.filter_bar.pack(fill=tk.X, padx=6, pady=(0, 2))

        columns = ("文件名", "类型", "大小", "修改时间")
        self.file_tree = ttk.Treeview(list_frame, columns=columns,
                                       show="tree headings", height=12,
                                       selectmode="extended")
        # #0 列：只有图标
        self.file_tree.heading("#0", text="")
        self.file_tree.column("#0", width=32, anchor="center", minwidth=32,
                              stretch=False)
        # 文件名列
        self.file_tree.heading("文件名", text="文件名",
                               command=lambda: self._sort_by_col("文件名"))
        self.file_tree.column("文件名", width=318, anchor="w", minwidth=60)
        # 其他列
        for col, text, w, anc in [("类型", "类型 ▼", 80, "center"),
                                    ("大小", "大小", 90, "center"),
                                    ("修改时间", "修改时间", 130, "center")]:
            self.file_tree.heading(col, text=text, command=lambda c=col: self._sort_by_col(c))
            self.file_tree.column(col, width=w, anchor=anc)

        # 行高和间距
        style = ttk.Style()
        style.configure("Treeview", rowheight=24)

        self.file_tree.bind("<Double-1>", self._on_file_double_click)
        self.file_tree.bind("<Button-3>", self._on_right_click)
        self.file_tree.bind("<<TreeviewSelect>>", self._on_file_select)

        self.context_menu = tk.Menu(self.root, tearoff=0, font=(FONT_FAMILY, 9))
        self.context_menu.add_command(label="打开", command=self._ctx_open)
        self.context_menu.add_command(label="打开所在文件夹", command=self._ctx_open_location)
        self.context_menu.add_separator()
        self.context_menu.add_command(label="移到回收站", command=self._ctx_trash)
        self.context_menu.add_command(label="永久删除", command=self._ctx_delete)

        vsb = ttk.Scrollbar(list_frame, orient="vertical", command=self.file_tree.yview)
        self.file_tree.configure(yscrollcommand=vsb.set)
        self.file_tree.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(6, 0), pady=4)
        vsb.pack(side=tk.RIGHT, fill=tk.Y, padx=(0, 6), pady=4)
        self.file_count_label = tk.Label(list_frame, text="", font=(FONT_FAMILY, 9),
                                          bg=COLOR_CARD, anchor=tk.W)
        self.file_count_label.pack(fill=tk.X, padx=8, pady=(0, 4))

        # 预览面板（可拖动分割）
        preview_frame = tk.Frame(paned, bg=COLOR_CARD, bd=1)
        paned.add(preview_frame, weight=1)

        tk.Label(preview_frame, text="预览", font=(FONT_FAMILY, 10, "bold"),
                 bg=COLOR_CARD).pack(pady=(4, 2), padx=8, anchor=tk.W)

        self.preview_canvas = tk.Canvas(preview_frame, height=120,
                                         highlightthickness=0)
        self.preview_canvas.pack(fill=tk.BOTH, padx=8, pady=(0, 6))

        self._preview_photo = None  # 防止图片被 GC

    # ── 排序 ──────────────────────────────────────────

    def _sort_by_col(self, col):
        if self._sort_col == col:
            self._sort_rev = not self._sort_rev
        else:
            self._sort_col, self._sort_rev = col, False
        self._apply_sort()

    def _apply_sort(self):
        if not self._file_data:
            return
        idx = {"文件名": 0, "类型": 1, "大小": 2, "修改时间": 3}[self._sort_col]
        def key(item):
            v = item[idx]
            if self._sort_col == "大小":
                try:
                    s = v.strip()
                    if "GB" in s: return float(s.replace("GB","").strip()) * 1024
                    elif "MB" in s: return float(s.replace("MB","").strip())
                    elif "KB" in s: return float(s.replace("KB","").strip()) / 1024
                    return 0
                except: return 0
            return v.lower() if isinstance(v, str) else v
        self._file_data.sort(key=key, reverse=self._sort_rev)
        # 更新表头排序箭头
        col_map = {"文件名": "文件名", "类型": "类型", "大小": "大小", "修改时间": "修改时间"}
        for c in ("文件名", "类型", "大小", "修改时间"):
            a = " ▲" if c == self._sort_col and self._sort_rev else (" ▼" if c == self._sort_col else "")
            self.file_tree.heading(col_map[c], text=(c if c != "类型" else "类型") + a)
        self._repopulate()

    def _rebuild_type_filters(self):
        """重建内置类型勾选框栏"""
        for w in self.filter_bar.winfo_children():
            w.destroy()
        types = sorted(set(d[1] for d in self._file_data))
        if not types:
            return
        # 标签
        tk.Label(self.filter_bar, text="类型:", bg=COLOR_CARD,
                 font=(FONT_FAMILY, 9)).pack(side=tk.LEFT, padx=(0, 4))
        self._filter_vars.clear()

        def apply_filter():
            selected = {t for t, var in self._filter_vars.items() if var.get()}
            if len(selected) == len(types):
                self._type_filter = set()
            else:
                self._type_filter = selected
            self._do_filter()

        for t in types:
            var = tk.BooleanVar(value=(not self._type_filter or t in self._type_filter))
            self._filter_vars[t] = var
            cb = tk.Checkbutton(self.filter_bar, text=t, variable=var,
                                bg=COLOR_CARD, font=(FONT_FAMILY, 9),
                                command=apply_filter)
            cb.pack(side=tk.LEFT, padx=2)

        # 全选/取消全选按钮
        def select_all():
            for var in self._filter_vars.values():
                var.set(True)
            apply_filter()

        def deselect_all():
            for var in self._filter_vars.values():
                var.set(False)
            apply_filter()

        ttk.Button(self.filter_bar, text="全选", command=select_all,
                   width=4).pack(side=tk.LEFT, padx=(4, 0))
        ttk.Button(self.filter_bar, text="无", command=deselect_all,
                   width=3).pack(side=tk.LEFT, padx=(2, 0))

    def _do_filter(self):
        """应用类型筛选并刷新列表"""
        self._repopulate()

    def _repopulate(self):
        self.file_tree.delete(*self.file_tree.get_children())
        self._file_path_map.clear()
        data = self._file_data if not self._type_filter else [d for d in self._file_data if d[1] in self._type_filter]
        for i, d in enumerate(data):
            tags = ["doc"] if d[5] else []
            icon = self._get_file_icon(d[4])
            kw = dict(iid=str(i), text="", values=(d[0], d[1], d[2], d[3]), tags=tags)
            if icon:
                kw['image'] = icon
            self.file_tree.insert("", "end", **kw)
            self._file_path_map[str(i)] = d[4]
        doc_count = sum(1 for d in data if d[5])
        self.file_count_label.config(text=f"共 {len(data)} 个文件（含 {doc_count} 个文档）")
        self.file_tree.tag_configure("doc", background=COLOR_DOC)
        self._rebuild_type_filters()

    def _get_file_icon(self, filepath):
        """从 Windows Shell 获取文件系统图标，按扩展名缓存"""
        ext = os.path.splitext(filepath)[1].lower()
        if ext in self._icon_cache:
            return self._icon_cache[ext]
        try:
            import win32gui
            import win32ui
            from PIL import Image, ImageTk
            shinfo = SHFILEINFOW()
            ret = ctypes.windll.shell32.SHGetFileInfoW(
                filepath, 0, ctypes.byref(shinfo), ctypes.sizeof(shinfo),
                SHGFI_ICON | SHGFI_SMALLICON | SHGFI_USEFILEATTRIBUTES
            )
            if not ret or not shinfo.hIcon:
                return None
            hicon = shinfo.hIcon
            icon_info = win32gui.GetIconInfo(hicon)
            hbm_color = icon_info[4]
            if not hbm_color:
                ctypes.windll.user32.DestroyIcon(hicon)
                return None
            bmp = win32ui.CreateBitmapFromHandle(hbm_color)
            info = bmp.GetInfo()
            w, h = info['bmWidth'], info['bmHeight']

            # 用 ctypes 的 GetDIBits 提取像素
            hdc_screen = win32gui.GetDC(0)
            hdc_mem = win32gui.CreateCompatibleDC(hdc_screen)
            hbm_old = win32gui.SelectObject(hdc_mem, hbm_color)

            class BHI_(ctypes.Structure):
                _fields_ = [
                    ('biSize', ctypes.c_uint32),
                    ('biWidth', ctypes.c_int32),
                    ('biHeight', ctypes.c_int32),
                    ('biPlanes', ctypes.c_uint16),
                    ('biBitCount', ctypes.c_uint16),
                    ('biCompression', ctypes.c_uint32),
                    ('biSizeImage', ctypes.c_uint32),
                    ('biXPelsPerMeter', ctypes.c_int32),
                    ('biYPelsPerMeter', ctypes.c_int32),
                    ('biClrUsed', ctypes.c_uint32),
                    ('biClrImportant', ctypes.c_uint32),
                ]

            class BMI_(ctypes.Structure):
                _fields_ = [('bmiHeader', BHI_)]

            bmi = BMI_()
            bmi.bmiHeader.biSize = ctypes.sizeof(BHI_)
            bmi.bmiHeader.biWidth = w
            bmi.bmiHeader.biHeight = -h
            bmi.bmiHeader.biPlanes = 1
            bmi.bmiHeader.biBitCount = 32
            bmi.bmiHeader.biCompression = 0

            pixels = (ctypes.c_ubyte * (w * h * 4))()
            gd = ctypes.windll.gdi32
            gd.GetDIBits.argtypes = [
                ctypes.c_void_p, ctypes.c_void_p,
                ctypes.c_uint, ctypes.c_uint,
                ctypes.c_void_p, ctypes.c_void_p, ctypes.c_uint,
            ]
            gd.GetDIBits.restype = ctypes.c_int

            gdi_ret = gd.GetDIBits(
                hdc_mem, int(hbm_color), 0, h,
                ctypes.cast(pixels, ctypes.c_void_p),
                ctypes.byref(bmi), 0,
            )

            photo = None
            if gdi_ret:
                img = Image.frombuffer('RGBA', (w, h), pixels, 'raw', 'BGRA', 0, 1)
                if w > 16 or h > 16:
                    img = img.resize((16, 16), Image.Resampling.LANCZOS)
                # 合成到白色背景，避免Tkinter透明通道问题
                bg = Image.new('RGBA', img.size, (255, 255, 255, 255))
                img = Image.alpha_composite(bg, img).convert('RGB')
                photo = ImageTk.PhotoImage(img)

            win32gui.SelectObject(hdc_mem, hbm_old)
            win32gui.DeleteDC(hdc_mem)
            win32gui.ReleaseDC(0, hdc_screen)
            ctypes.windll.user32.DestroyIcon(hicon)
            self._icon_cache[ext] = photo
            return photo
        except Exception:
            self._icon_cache[ext] = None
            return None

    def _debounce_search(self):
        """搜索防抖：停止打字 300ms 后才触发搜索"""
        if hasattr(self, '_search_after_id'):
            self.root.after_cancel(self._search_after_id)
        self._search_after_id = self.root.after(300, self._refresh_file_list)

    # ── 文件列表 ───────────────────────────────────────

    def _refresh_file_list(self):
        if not self.current_dir or not os.path.isdir(self.current_dir):
            return
        kw = self.search_entry.get().strip().lower()
        def collect(base):
            r = []
            bn = os.path.normpath(base)
            try:
                with os.scandir(base) as it:
                    entries = sorted(it, key=lambda e: e.name.lower())
                    for entry in entries:
                        if len(r) >= 500:
                            break
                        if entry.is_dir(follow_symlinks=False):
                            name = entry.name.lower()
                            if name.startswith(".") or name.startswith("$") or name in SKIP_DIRS_LOWER:
                                continue
                            # 递归进入子目录
                            sub_files = collect_sub(entry.path, 4)
                            r.extend(sub_files)
                            if len(r) >= 500:
                                break
                        elif entry.is_file():
                            fn = entry.name
                            if not should_show_file(fn):
                                continue
                            if kw and kw not in fn.lower():
                                continue
                            fp = entry.path
                            ext = os.path.splitext(fn)[1].lower()
                            try:
                                st = entry.stat(follow_symlinks=False)
                                sz, tm = format_size(st.st_size), datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M")
                            except: sz, tm = "?", "?"
                            r.append((fn, get_type_label(ext), sz, tm, fp, ext in (".doc", ".docx", ".wps")))
            except: pass
            return r

        def collect_sub(base, depth):
            """递归收集子目录"""
            r = []
            if depth <= 0:
                return r
            try:
                with os.scandir(base) as it:
                    entries = sorted(it, key=lambda e: e.name.lower())
                    for entry in entries:
                        if len(r) >= 500:
                            break
                        if entry.is_dir(follow_symlinks=False):
                            name = entry.name.lower()
                            if name.startswith(".") or name.startswith("$") or name in SKIP_DIRS_LOWER:
                                continue
                            r.extend(collect_sub(entry.path, depth - 1))
                            if len(r) >= 500:
                                break
                        elif entry.is_file():
                            fn = entry.name
                            if not should_show_file(fn):
                                continue
                            if kw and kw not in fn.lower():
                                continue
                            fp = entry.path
                            ext = os.path.splitext(fn)[1].lower()
                            try:
                                st = entry.stat(follow_symlinks=False)
                                sz, tm = format_size(st.st_size), datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M")
                            except: sz, tm = "?", "?"
                            r.append((fn, get_type_label(ext), sz, tm, fp, ext in (".doc", ".docx", ".wps")))
            except: pass
            return r
        threading.Thread(target=lambda: (
            setattr(self, '_file_data', collect(self.current_dir)),
            self.root.after(0, lambda: (setattr(self, '_sort_col', None),
                setattr(self, '_sort_rev', False), setattr(self, '_type_filter', set()),
                self.file_tree.heading("类型", text="类型 ▼"), self._repopulate()))
        ), daemon=True).start()

    def _on_file_double_click(self, event):
        item = self.file_tree.focus()
        if item:
            fp = self._file_path_map.get(item)
            if fp and os.path.isfile(fp):
                try:
                    os.startfile(fp)
                    self._add_recent(fp)
                except Exception as e:
                    messagebox.showerror("打开失败", str(e))

    def _on_file_select(self, event):
        item = self.file_tree.focus()
        if item:
            fp = self._file_path_map.get(item)
            if fp:
                # 显示加载提示
                self.preview_canvas.delete("all")
                self.preview_canvas.create_text(200, 40, text="加载中...",
                                                  fill="#999", font=(FONT_FAMILY, 9))
                # 防抖 + 后台线程
                if hasattr(self, '_preview_after_id'):
                    self.root.after_cancel(self._preview_after_id)
                self._preview_after_id = self.root.after(200, lambda p=fp: self._preview_thread(p))

    def _preview_thread(self, filepath):
        """后台线程预览，避免卡顿"""
        def worker():
            try:
                # 先在后台读取文件内容
                ext = os.path.splitext(filepath)[1].lower()
                info_text = self._extract_preview_data(filepath, ext)
                self.root.after(0, lambda: self._draw_preview(filepath, ext, info_text))
            except Exception as e:
                self.root.after(0, lambda: self._draw_preview(filepath, "", f"预览失败: {e}"))
        threading.Thread(target=worker, daemon=True).start()

    def _extract_preview_data(self, filepath, ext):
        """后台线程提取预览内容"""
        name = os.path.basename(filepath)
        info = f"文件: {name}"
        try:
            st = os.stat(filepath)
            info += f"  |  {format_size(st.st_size)}  |  {datetime.fromtimestamp(st.st_mtime).strftime('%Y-%m-%d %H:%M')}"
        except: pass

        if ext in {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff"}:
            img = Image.open(filepath)
            img.thumbnail((500, 250), Image.Resampling.LANCZOS)
            return ("image", info, img)
        elif ext in TEXT_EXTENSIONS:
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                lines = f.readlines()[:5]
            return ("text", info, "".join(lines)[:300])
        elif ext in {".doc", ".wps", ".rtf", ".odt"}:
            import win32com.client as w32
            import pythoncom
            pythoncom.CoInitialize()
            word = w32.Dispatch("Word.Application")
            word.Visible = False
            word.DisplayAlerts = 0
            doc = word.Documents.Open(os.path.abspath(filepath), ReadOnly=True, AddToRecentFiles=False)
            text = doc.Content.Text[:500]
            doc.Close(SaveChanges=0)
            word.Quit()
            pythoncom.CoUninitialize()
            return ("text", info, text.strip() or "(空白文档)")
        elif ext == ".docx":
            from docx import Document
            doc = Document(filepath)
            paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()][:6]
            return ("text", info, "\n".join(paras)[:400] or "(空白文档)")
        elif ext in {".ppt", ".pptx"}:
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = []
            for slide in prs.slides[:3]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text.strip())
            return ("text", info, "\n".join(texts)[:400] or "(空白演示)")
        elif ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(filepath)
            text = ""
            if len(reader.pages) > 0:
                text = reader.pages[0].extract_text()[:400]
            return ("text", info, text or "(无法提取文本)")
        elif ext == ".csv":
            import csv
            with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
                reader = csv.reader(f)
                rows = [",".join(row[:4]) for row in list(reader)[:5]]
            return ("text", info, "\n".join(rows)[:400])
        else:
            return ("text", info, f"[{get_type_label(ext)}文件] 双击打开")

    def _draw_preview(self, filepath, ext, preview_data):
        """主线程绘制预览"""
        c = self.preview_canvas
        c.delete("all")
        self._preview_photo = None

        if preview_data is None:
            return

        ptype = preview_data[0]
        info = preview_data[1]

        c.create_text(10, 10, text=info, anchor="nw", font=(FONT_FAMILY, 9), fill="#333")

        if ptype == "image":
            img = preview_data[2]
            cw = max(c.winfo_width(), 400)
            ch = max(c.winfo_height(), 120)
            self._preview_photo = ImageTk.PhotoImage(img)
            c.create_image(cw // 2, ch // 2 + 10, image=self._preview_photo, anchor="center")
        else:
            text = preview_data[2]
            c.create_text(10, 30, text=text, anchor="nw",
                          font=(FONT_FAMILY, 9), fill="#555",
                          width=max(c.winfo_width() - 20, 380))

    def _on_right_click(self, event):
        item = self.file_tree.identify_row(event.y)
        if item:
            self.file_tree.selection_set(item)
            self.file_tree.focus(item)
            self.context_menu.post(event.x_root, event.y_root)

        item = self.file_tree.focus()
        return self._file_path_map.get(item) if item else None

    def _get_selected_path(self):
        item = self.file_tree.focus()
        return self._file_path_map.get(item) if item else None

    def _ctx_open(self):
        fp = self._get_selected_path()
        if fp and os.path.isfile(fp):
            try:
                os.startfile(fp)
                self._add_recent(fp)
            except Exception as e: messagebox.showerror("打开失败", str(e))

    def _ctx_open_location(self):
        fp = self._get_selected_path()
        if fp and os.path.isfile(fp):
            try: os.startfile(os.path.dirname(fp))
            except Exception as e: messagebox.showerror("打开失败", str(e))

    def _ctx_trash(self):
        fp = self._get_selected_path()
        if fp and os.path.isfile(fp) and messagebox.askyesno("确认", f"移到回收站？\n{os.path.basename(fp)}"):
            try:
                send2trash(fp)
                self.status_var = getattr(self, 'status_var', None) or tk.StringVar(value="")
                self._refresh_file_list()
            except Exception as e: messagebox.showerror("失败", str(e))

    def _ctx_delete(self):
        fp = self._get_selected_path()
        if fp and os.path.isfile(fp) and messagebox.askyesno("永久删除", f"不可恢复！\n{os.path.basename(fp)}", icon="warning"):
            try:
                os.remove(fp)
                self._refresh_file_list()
            except Exception as e: messagebox.showerror("删除失败", str(e))

    # ── 文件夹管理 ─────────────────────────────────────

    def _data_path(self):
        return os.path.join(os.path.dirname(os.path.abspath(sys.argv[0])), "docuhub_folders.json")

    def _save_folders(self):
        try:
            with open(self._data_path(), "w", encoding="utf-8") as f:
                json.dump({"locked": list(self.locked_folders), "folders": self.saved_folders}, f, ensure_ascii=False, indent=2)
        except: pass

    def _load_folders(self):
        p = self._data_path()
        if not os.path.isfile(p): return
        try:
            with open(p, "r", encoding="utf-8") as f: d = json.load(f)
            self.locked_folders = set(d.get("locked", []))
            self.saved_folders = [x for x in d.get("folders", []) if os.path.isdir(x)]
            for fp in self.saved_folders:
                n = os.path.basename(os.path.normpath(fp))
                tag = "🔒" if fp in self.locked_folders else "📂"
                self.folder_list.append((f"{tag} {n}", fp))
                self.folder_listbox.insert(tk.END, self.folder_list[-1][0])
            if self.folder_list:
                self.folder_listbox.selection_set(0)
                self._on_folder_select(None)
        except: pass

    def _add_folder(self):
        path = filedialog.askdirectory(title="选择文件夹")
        if not path: return
        n = os.path.basename(os.path.normpath(path))
        self.folder_list.append((f"🔒 {n}", path))
        self.folder_listbox.insert(tk.END, f"🔒 {n}")
        self.saved_folders.append(path)
        self.locked_folders.add(path)
        self._save_folders()
        self.folder_listbox.selection_set(len(self.folder_list) - 1)
        self._on_folder_select(None)

    def _start_full_scan(self):
        self.folder_listbox.delete(0, tk.END)
        self.folder_listbox.insert(tk.END, "扫描中...")
        self.scan_btn.config(state=tk.DISABLED)
        def scan():
            import string
            results, seen = [], set()
            for l in string.ascii_uppercase:
                d = f"{l}:\\"
                if not os.path.exists(d): continue
                try:
                    for root, dirs, files in os.walk(d, topdown=True):
                        dirs[:] = [x for x in dirs if x.lower() not in SKIP_DIRS and not x.startswith("$")]
                        p = os.path.normpath(root)
                        if p not in seen and any(os.path.isfile(os.path.join(root, f)) and should_show_file(f) for f in files):
                            results.append((f"📁 {os.path.basename(root)}", root))
                            seen.add(p)
                        if root.replace(d, "").count(os.sep) >= 4: dirs.clear()
                except: pass
                if len(results) > 500: break
            results.sort(key=lambda x: x[1])
            self.root.after(0, self._on_scan_done, results)
        threading.Thread(target=scan, daemon=True).start()

    def _on_scan_done(self, folders):
        self.folder_listbox.delete(0, tk.END)
        self.folder_list = folders
        self.scan_btn.config(state=tk.NORMAL)
        for l, _ in folders: self.folder_listbox.insert(tk.END, l)
        if folders:
            self.folder_listbox.selection_set(0)
            self._on_folder_select(None)

    def _on_folder_select(self, event):
        sel = self.folder_listbox.curselection()
        if not sel or sel[0] >= len(self.folder_list): return
        path = self.folder_list[sel[0]][1]
        if not os.path.isdir(path): return
        self.current_dir = path
        self.dir_label.config(text=path)
        self._refresh_file_list()

    def _on_folder_right_click(self, event):
        idx = self.folder_listbox.nearest(event.y)
        if 0 <= idx < len(self.folder_list):
            self.folder_listbox.selection_clear(0, tk.END)
            self.folder_listbox.selection_set(idx)
            self.folder_menu.post(event.x_root, event.y_root)

    def _folder_refresh(self): self._on_folder_select(None)
    def _folder_explorer(self):
        sel = self.folder_listbox.curselection()
        if sel and sel[0] < len(self.folder_list):
            try: os.startfile(self.folder_list[sel[0]][1])
            except: pass
    def _folder_lock(self):
        sel = self.folder_listbox.curselection()
        if not sel or sel[0] >= len(self.folder_list): return
        idx, path = sel[0], self.folder_list[sel[0]][1]
        if path in self.locked_folders: return
        self.locked_folders.add(path)
        if path not in self.saved_folders: self.saved_folders.append(path)
        self._save_folders()
        n = os.path.basename(os.path.normpath(path))
        self.folder_list[idx] = (f"🔒 {n}", path)
        self.folder_listbox.delete(idx); self.folder_listbox.insert(idx, f"🔒 {n}")
        self.folder_listbox.selection_set(idx)
    def _folder_unlock(self):
        sel = self.folder_listbox.curselection()
        if not sel or sel[0] >= len(self.folder_list): return
        idx, path = sel[0], self.folder_list[sel[0]][1]
        self.locked_folders.discard(path); self._save_folders()
        n = os.path.basename(os.path.normpath(path))
        self.folder_list[idx] = (f"📂 {n}", path)
        self.folder_listbox.delete(idx); self.folder_listbox.insert(idx, f"📂 {n}")
        self.folder_listbox.selection_set(idx)
    def _folder_remove(self):
        sel = self.folder_listbox.curselection()
        if not sel or sel[0] >= len(self.folder_list): return
        idx, path = sel[0], self.folder_list[sel[0]][1]
        if path in self.locked_folders:
            messagebox.showwarning("禁止", "已锁定，请先解锁"); return
        if not messagebox.askyesno("确认", "从列表删除？"): return
        self.folder_list.pop(idx); self.folder_listbox.delete(idx)
        self.saved_folders = [p for p in self.saved_folders if p != path]
        self.locked_folders.discard(path); self._save_folders()
        if self.folder_list: self.folder_listbox.selection_set(min(idx, len(self.folder_list)-1))

    # ── 最近打开 ───────────────────────────────────────

    def _load_recent(self):
        if os.path.isfile(self.RECENT_PATH):
            try:
                with open(self.RECENT_PATH, "r", encoding="utf-8") as f:
                    self.recent_files = json.load(f)
            except: self.recent_files = []
        self._refresh_recent_list()

    def _save_recent(self):
        try:
            with open(self.RECENT_PATH, "w", encoding="utf-8") as f:
                json.dump(self.recent_files, f, ensure_ascii=False, indent=2)
        except: pass

    def _add_recent(self, filepath):
        filepath = os.path.normpath(filepath)
        self.recent_files = [p for p in self.recent_files if p != filepath]
        self.recent_files.insert(0, filepath)
        self.recent_files = self.recent_files[:30]
        self._save_recent()
        self._refresh_recent_list()

    def _refresh_recent_list(self):
        self.recent_listbox.delete(0, tk.END)
        for fp in self.recent_files[:10]:
            n = os.path.basename(fp)
            self.recent_listbox.insert(tk.END, n)

    def _open_recent(self, event):
        sel = self.recent_listbox.curselection()
        if not sel: return
        idx = sel[0]
        if idx < len(self.recent_files):
            fp = self.recent_files[idx]
            if os.path.isfile(fp):
                try:
                    os.startfile(fp)
                    self._add_recent(fp)
                except Exception as e:
                    messagebox.showerror("打开失败", str(e))
            else:
                messagebox.showinfo("文件不存在", f"{fp}\n可能已被移动或删除")

    def _recent_right_click(self, event):
        idx = self.recent_listbox.nearest(event.y)
        if 0 <= idx < len(self.recent_files):
            self.recent_listbox.selection_clear(0, tk.END)
            self.recent_listbox.selection_set(idx)
            self.recent_menu.post(event.x_root, event.y_root)

    def _recent_open(self):
        sel = self.recent_listbox.curselection()
        if not sel: return
        idx = sel[0]
        if idx < len(self.recent_files):
            fp = self.recent_files[idx]
            if os.path.isfile(fp):
                try:
                    os.startfile(fp)
                    self._add_recent(fp)
                except Exception as e:
                    messagebox.showerror("打开失败", str(e))
            else:
                messagebox.showinfo("文件不存在", f"{fp}\n可能已被移动或删除")

    def _recent_open_location(self):
        sel = self.recent_listbox.curselection()
        if not sel: return
        idx = sel[0]
        if idx < len(self.recent_files):
            fp = self.recent_files[idx]
            if os.path.isfile(fp):
                try: os.startfile(os.path.dirname(fp))
                except Exception as e: messagebox.showerror("打开失败", str(e))
            else:
                messagebox.showinfo("文件不存在", f"{fp}\n可能已被移动或删除")

    def _recent_remove(self):
        sel = self.recent_listbox.curselection()
        if not sel: return
        idx = sel[0]
        if idx < len(self.recent_files):
            self.recent_files.pop(idx)
            self._save_recent()
            self._refresh_recent_list()

    # ── 全盘搜索 ───────────────────────────────────────

    def _start_global_search(self):
        kw = self.global_search_var.get().strip()
        if not kw: return
        self.dir_label.config(text=f"搜索：{kw}", fg=COLOR_ACCENT)
        self.file_count_label.config(text="搜索中...")
        def search():
            import string
            results = []
            for l in string.ascii_uppercase:
                d = f"{l}:\\"
                if not os.path.exists(d): continue
                try:
                    for root, dirs, files in os.walk(d, topdown=True):
                        dirs[:] = [x for x in dirs if x.lower() not in SKIP_DIRS and not x.startswith("$")]
                        for fn in files:
                            if kw.lower() in fn.lower() and should_show_file(fn):
                                fp = os.path.join(root, fn)
                                try:
                                    st = os.stat(fp)
                                    results.append((fn, get_type_label(os.path.splitext(fn)[1]),
                                                    format_size(st.st_size),
                                                    datetime.fromtimestamp(st.st_mtime).strftime("%Y-%m-%d %H:%M"),
                                                    fp, os.path.splitext(fn)[1].lower() in (".doc", ".docx")))
                                except: pass
                        if root.replace(d, "").count(os.sep) >= 4: dirs.clear()
                except: pass
                if len(results) > 500: break
            self.root.after(0, self._on_search_done, results)
        threading.Thread(target=search, daemon=True).start()

    def _on_search_done(self, results):
        self._file_data = results
        self._sort_col = self._sort_rev = None
        self._type_filter = set()
        self.file_tree.heading("类型", text="类型 ▼")
        self._repopulate()


if __name__ == "__main__":
    DocuHubApp().run()
